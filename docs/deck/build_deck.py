"""Fill the provided SIH template with Eco-Loop content.

Template rules honoured: six slides including the title page, the instructions slide removed,
section headings kept as given, points and visuals rather than paragraphs.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
FIG = HERE / "fig"

INK = RGBColor(0x1F, 0x38, 0x64)
NAVY = RGBColor(0x1F, 0x49, 0x7D)
BLUE = RGBColor(0x4F, 0x81, 0xBD)
GOOD = RGBColor(0x2E, 0x7D, 0x32)
BAD = RGBColor(0xC0, 0x50, 0x4D)
BODY = RGBColor(0x33, 0x41, 0x5C)
MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xEE, 0xF2, 0xF7)

FONT = "Calibri"


def drop_slide(prs, index):
    slides = prs.slides._sldIdLst
    rid = slides[index].rId
    prs.part.drop_rel(rid)
    slides.remove(slides[index])


def clear(shape):
    shape._element.getparent().remove(shape._element)


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def line(frame, text, size, colour=BODY, bold=False, italic=False, space_before=0,
         space_after=4, first=False, align=PP_ALIGN.LEFT):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = FONT
    return para


def card(slide, left, top, width, height, fill=CARD):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.06
    if shape.has_text_frame:
        shape.text_frame.clear()
    return shape


def stat(slide, left, top, width, value, label, colour):
    frame = textbox(slide, left, top, width, 0.95)
    line(frame, value, 30, colour, bold=True, first=True, space_after=0, align=PP_ALIGN.CENTER)
    line(frame, label, 11, MUTED, space_before=2, align=PP_ALIGN.CENTER)


prs = Presentation(str(HERE / "template.pptx"))
drop_slide(prs, 0)  # instructions slide
slides = list(prs.slides)

# renumber the footer page numbers left behind by the deletion
for number, slide in enumerate(slides, start=1):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip().isdigit():
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = str(number)
                    break
                break

# ----------------------------------------------------------------- 1. title
s = slides[0]
for shape in list(s.shapes):
    if shape.has_text_frame and "Problem Statement ID" in shape.text_frame.text:
        clear(shape)
    if shape.has_text_frame and "TITLE PAGE" in shape.text_frame.text:
        clear(shape)

f = textbox(s, 0.7, 0.55, 12.0, 1.5)
line(f, "Eco-Loop Building Agents", 40, INK, bold=True, first=True, space_after=6)
line(f, "A supervisory controller that commissions itself onto a building it has never seen —"
        " and refuses to deploy when nothing earns its place.", 15, MUTED)

f = textbox(s, 0.7, 2.35, 6.0, 3.0)
for label, value in [
    ("Problem Statement ID", "«enter your PS ID»"),
    ("Problem Statement Title", "Eco-Loop Building Agents"),
    ("Theme", "«enter theme»"),
    ("PS Category", "Software"),
    ("Team / Student Name", "«enter name as registered»"),
    ("Student ID", "«enter ID»"),
]:
    para = f.paragraphs[0] if label == "Problem Statement ID" else f.add_paragraph()
    para.space_after = Pt(9)
    a = para.add_run()
    a.text = f"{label} — "
    a.font.size, a.font.bold, a.font.color.rgb, a.font.name = Pt(14), True, INK, FONT
    b = para.add_run()
    b.text = value
    b.font.size, b.font.color.rgb, b.font.name = Pt(14), BODY, FONT

card(s, 7.05, 2.35, 5.6, 2.30, CARD)
f = textbox(s, 7.4, 2.62, 4.9, 2.2)
line(f, "FULL YEAR, DOE MEDIUM OFFICE, CHICAGO", 10, MUTED, bold=True, first=True, space_after=10)
for value, label, colour in [
    ("−5.33%", "electricity", GOOD),
    ("−8.66%", "peak demand", GOOD),
    ("56.0 → 51.5", "unmet hours (lower than baseline)", BLUE),
]:
    para = f.add_paragraph()
    para.space_after = Pt(7)
    a = para.add_run()
    a.text = value + "   "
    a.font.size, a.font.bold, a.font.color.rgb, a.font.name = Pt(19), True, colour, FONT
    b = para.add_run()
    b.text = label
    b.font.size, b.font.color.rgb, b.font.name = Pt(12), BODY, FONT

f = textbox(s, 7.4, 4.78, 5.3, 0.4)
line(f, "every figure regenerated by one command:  make evidence", 11, MUTED,
     italic=True, first=True)

card(s, 0.65, 5.55, 12.0, 1.05, RGBColor(0xF4, 0xF6, 0xF8))
f = textbox(s, 0.95, 5.78, 11.4, 0.7)
line(f, "Closed loop on a live EnergyPlus instance  ·  13-tool MCP server, stdio and HTTP  ·  "
        "commissions itself onto a building it has never seen  ·  repairs a model EnergyPlus "
        "refuses  ·  55 tests", 12.5, BODY, first=True, align=PP_ALIGN.CENTER)

# ------------------------------------------------------------- 2. the idea
s = slides[1]
for shape in list(s.shapes):
    if shape.has_text_frame and "Proposed Solution" in shape.text_frame.text:
        clear(shape)

f = textbox(s, 0.7, 1.30, 12.0, 0.45)
line(f, "Eco-Loop — a supervisory controller that commissions itself, and refuses when nothing "
        "earns its place", 18, INK, bold=True, first=True, space_after=5)

f = textbox(s, 0.7, 1.82, 12.0, 0.5)
line(f, "Advanced control already saves 15–30%. It reaches under 5% of buildings, because every "
        "deployment costs weeks of engineering — the blocker is commissioning, not the algorithm.",
     13.5, BODY, italic=True, first=True)

card(s, 0.65, 2.50, 6.05, 2.90)
f = textbox(s, 0.95, 2.72, 5.45, 2.6)
line(f, "WHAT IT DOES", 11, NAVY, bold=True, first=True, space_after=9)
for step, text in [
    ("1  Discover", "reads the model's own wiring to find what can be actuated"),
    ("2  Survey", "runs the building untouched to see where the energy goes"),
    ("3  Trial", "simulates each candidate measure against that baseline"),
    ("4  Deploy or refuse", "keeps only what actually helps; ships nothing otherwise"),
]:
    para = f.add_paragraph()
    para.space_after = Pt(8)
    a = para.add_run()
    a.text = step + " — "
    a.font.size, a.font.bold, a.font.color.rgb, a.font.name = Pt(13), True, INK, FONT
    b = para.add_run()
    b.text = text
    b.font.size, b.font.color.rgb, b.font.name = Pt(12.5), BODY, FONT
line(f, "Control is written into a live EnergyPlus instance every timestep through its runtime "
        "API — not by editing a file and re-running.", 12, MUTED, italic=True, space_before=6)

card(s, 6.95, 2.50, 5.7, 2.90, RGBColor(0xF7, 0xF0, 0xEE))
f = textbox(s, 7.25, 2.72, 5.1, 2.6)
line(f, "WHY IT IS DIFFERENT", 11, BAD, bold=True, first=True, space_after=9)
for head, text in [
    ("A measure must earn its place",
     "applicability is not benefit. On one building the obvious measure targeted 16.6% of "
     "load, was fully actuable, and made things 2.14% worse. Only a trial caught it."),
    ("It refuses",
     "two of three building-and-climate pairs deploy nothing, and say why."),
    ("Comfort is scored where the controller cannot reach",
     "unmet hours are gameable — a zone held at 29 °C reports zero. A test proves it."),
    ("We can prove where the LLM helps",
     "a perfect-foresight arm bounds the task, so a null result is visible instead of hidden."),
]:
    para = f.add_paragraph()
    para.space_after = Pt(7)
    a = para.add_run()
    a.text = head + "  "
    a.font.size, a.font.bold, a.font.color.rgb, a.font.name = Pt(12), True, INK, FONT
    b = para.add_run()
    b.text = text
    b.font.size, b.font.color.rgb, b.font.name = Pt(11), BODY, FONT

for i, (value, label, colour) in enumerate([
    ("−5.33%", "electricity, full year", GOOD),
    ("−8.66%", "peak demand", GOOD),
    ("56.0 → 51.5", "unmet hours, improved", BLUE),
    ("0", "measures shipped where none helped", BAD),
]):
    stat(s, 0.65 + i * 3.06, 5.60, 2.9, value, label, colour)

# ------------------------------------------------------- 3. technical approach
s = slides[2]
for shape in list(s.shapes):
    if shape.has_text_frame and "Technologies to be used" in shape.text_frame.text:
        clear(shape)

s.shapes.add_picture(str(FIG / "architecture.png"), Inches(0.62), Inches(1.15), width=Inches(12.1))

f = textbox(s, 0.66, 5.95, 12.0, 1.0)
line(f, "Python 3.11 · EnergyPlus 26.1 runtime API (pyenergyplus) · epJSON model mutation · "
        "Qwen2.5-3B on Ollama (local, open-source) · Model Context Protocol (FastMCP) · "
        "Pydantic contracts · pytest, 55 tests · matplotlib report",
     11.5, BODY, first=True, space_after=6)
line(f, "The model is never in the inner loop: 52,560 inferences a year at 7.6 s each is 111 hours "
        "against a 14-second simulation. It sets one number a day; the deterministic core runs "
        "every timestep and the Guardian bounds it.", 11.5, MUTED, italic=True)

# --------------------------------------------------- 4. feasibility and viability
s = slides[3]
for shape in list(s.shapes):
    if shape.has_text_frame and "Analysis of the feasibility" in shape.text_frame.text:
        clear(shape)

card(s, 0.65, 1.15, 5.6, 2.35)
f = textbox(s, 0.95, 1.4, 5.0, 2.0)
line(f, "BUILT AND RUNNING TODAY", 11, GOOD, bold=True, first=True, space_after=8)
for text in [
    "Whole system runs on one laptop — an M1 with 5.3 GiB of VRAM",
    "Annual simulation with full telemetry: 14 seconds",
    "55 tests, including real simulations, green in CI",
    "One command regenerates every published number in ~3 minutes",
]:
    line(f, "•  " + text, 12, BODY, space_after=6)

card(s, 6.55, 1.15, 6.1, 2.35, RGBColor(0xF7, 0xF0, 0xEE))
f = textbox(s, 6.85, 1.4, 5.5, 2.0)
line(f, "RISKS, AND WHAT ANSWERS THEM", 11, BAD, bold=True, first=True, space_after=8)
for risk, fix in [
    ("EnergyPlus API fails silently", "four known traps handled; totals reconciled against E+'s own report by test"),
    ("Model slow, wrong or offline", "never in the inner loop; last-good plan; a test kills it mid-run"),
    ("A measure harms the building", "trial verification, then refusal"),
    ("Comfort quietly traded away", "fixed band the controller cannot move, enforced by the Guardian"),
]:
    para = f.add_paragraph()
    para.space_after = Pt(5)
    a = para.add_run()
    a.text = risk + " → "
    a.font.size, a.font.bold, a.font.color.rgb, a.font.name = Pt(11.5), True, INK, FONT
    b = para.add_run()
    b.text = fix
    b.font.size, b.font.color.rgb, b.font.name = Pt(11), BODY, FONT

card(s, 0.65, 3.75, 12.0, 2.15, RGBColor(0xEC, 0xF1, 0xE9))
f = textbox(s, 0.95, 4.0, 11.4, 2.0)
line(f, "THE RESULT WE DID NOT EXPECT, AND KEPT", 11, GOOD, bold=True, first=True, space_after=8)
line(f, "The LLM-supervised arm does not beat the deterministic controller: −1.50% against −5.33% "
        "over a year.", 13, BODY, space_after=6)
line(f, "Before blaming the model we bounded the task. A perfect-foresight arm — the same "
        "controller handed the day's forecast peak, which is exactly the anticipation the model "
        "was asked for — lands on the reactive controller to within 0.03 points over a season. "
        "There is no headroom on this measure for anything to find.", 12.5, BODY, space_after=6)
line(f, "Against the untouched building alone the agent saves 1.5% and could be reported as a "
        "success. It is only visible as a null result because the deterministic bar exists. The "
        "model's value here is commissioning and refusal, not energy — and the catalogue of "
        "measures, not the framework, is now the limit.", 12.5, MUTED, italic=True)

# ---------------------------------------------------------------- 5. artifacts
s = slides[4]
for shape in list(s.shapes):
    if shape.has_text_frame and "Relevant artifacts" in shape.text_frame.text:
        clear(shape)

f = textbox(s, 0.66, 1.10, 12.0, 0.3)
line(f, "ABLATION LADDER \u2014 five arms, identical weather, run period and timestep, compared on "
        "one asserted-identical clock", 11, NAVY, bold=True, first=True)
s.shapes.add_picture(str(FIG / "ladder_wide.png"), Inches(0.66), Inches(1.44), width=Inches(11.9))

f = textbox(s, 0.66, 4.10, 12.0, 0.3)
line(f, "SELF-COMMISSIONING \u2014 the same code pointed at a building it was not built for",
     11, NAVY, bold=True, first=True)
s.shapes.add_picture(str(FIG / "terminal.png"), Inches(0.66), Inches(4.42), width=Inches(11.9))

f = textbox(s, 0.66, 6.58, 12.0, 0.45)
line(f, "Repository: \u00abpaste your GitHub URL\u00bb   \u00b7   source, 9 architecture decision "
        "records, the generated evidence report, and the agent\u2019s 390-record decision journal "
        "so the LLM arm replays exactly without a model server", 10.5, MUTED, first=True)

# --------------------------------------------------------------- 6. references
s = slides[5]
for shape in list(s.shapes):
    if shape.has_text_frame and "Details / Links" in shape.text_frame.text:
        clear(shape)

left = [
    ("ASHRAE Guideline 36-2021", "High-performance sequences of operation for HVAC systems — "
     "the supply-air reset and trim-and-respond logic implemented here"),
    ("ASHRAE 90.1 Appendix G", "Unmet load hours as a compliance limit; our baseline sits at "
     "56 h against the 300 h allowance"),
    ("DOE Commercial Reference Buildings", "Medium and small office prototypes used as the "
     "baseline models, unmodified"),
]
right = [
    ("EnergyPlus 26.1 Python API", "nrel.github.io/EnergyPlus — runtime callbacks, actuators, "
     "and the epJSON schema used for model mutation"),
    ("Model Context Protocol", "modelcontextprotocol.io — tool specification; served here over "
     "stdio and streamable HTTP via FastMCP"),
    ("Qwen2.5-3B-Instruct · Ollama", "open-source model run locally with JSON-schema-constrained "
     "decoding; ISHRAE 2014 New Delhi weather from climate.onebuilding.org"),
]

for column, items in ((0.65, left), (6.95, right)):
    f = textbox(s, column, 1.3, 5.7, 4.6)
    firstline = True
    for title, detail in items:
        para = f.paragraphs[0] if firstline else f.add_paragraph()
        para.space_after = Pt(3)
        para.space_before = Pt(0 if firstline else 14)
        run = para.add_run()
        run.text = title
        run.font.size, run.font.bold, run.font.color.rgb, run.font.name = Pt(14), True, INK, FONT
        line(f, detail, 12, BODY, space_after=0)
        firstline = False

card(s, 0.65, 4.45, 12.0, 0.92, RGBColor(0xF4, 0xF6, 0xF8))
f = textbox(s, 0.95, 4.66, 11.4, 0.75)
line(f, "VERIFY IT YOURSELF", 10.5, NAVY, bold=True, first=True, space_after=6)
line(f, "make setup   \u00b7   make evidence   \u2014  regenerates every published number in about "
        "three minutes, with no model server needed   \u00b7   make test  \u2014  55 tests", 12,
     BODY, space_after=0)

f = textbox(s, 0.65, 5.72, 12.0, 0.6)
line(f, "Method note — forecasts given to the model are deliberately degraded with error growing "
        "by lead time, because the weather file is perfect foresight. Delhi results are a "
        "climate-sensitivity probe, not a claim about Indian buildings: the US prototype is "
        "undersized for that climate.", 11, MUTED, italic=True, first=True)

out = HERE / "Eco-Loop Building Agents - Eco-Loop.pptx"
prs.save(str(out))
print("wrote", out)
